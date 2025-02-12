from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_shape = slide.shapes.title
    subtitle_shape = slide.placeholders[1]
    
    title_shape.text = title
    subtitle_shape.text = subtitle
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)  # Green color
    
    return slide

def add_content_slide(prs, title, content):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    content_shape = slide.placeholders[1]
    
    title_shape.text = title
    content_frame = content_shape.text_frame
    
    for line in content:
        p = content_frame.add_paragraph()
        p.text = line
        p.font.size = Pt(18)
    
    # Style the title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.color.rgb = RGBColor(46, 125, 50)
    
    return slide

def create_presentation():
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, 
                   "Bienes Raíces con IA",
                   "Revolucionando el Sector Inmobiliario con Inteligencia Artificial")
    
    # Introduction
    add_content_slide(prs,
                     "Introducción",
                     ["• Plataforma integral de gestión inmobiliaria",
                      "• Impulsada por Inteligencia Artificial",
                      "• Asistente virtual 24/7",
                      "• Búsqueda inteligente de propiedades",
                      "• Sistema de gestión de contenidos"])
    
    # Features
    add_content_slide(prs,
                     "Características Principales",
                     ["Portal Web Inteligente:",
                      "• Búsqueda avanzada de propiedades",
                      "• Filtros personalizables",
                      "• Blog integrado",
                      "• Asistente virtual de IA"])
    
    # Benefits
    add_content_slide(prs,
                     "Beneficios para Apeko",
                     ["Optimización Operativa:",
                      "• Reducción de tiempo en tareas administrativas",
                      "• Automatización de procesos",
                      "• Mejor gestión de leads",
                      "• Análisis de datos en tiempo real"])
    
    # Licensing
    add_content_slide(prs,
                     "Modelo de Licenciamiento",
                     ["Tipos de Licencias:",
                      "• Básica: $499/mes - 100 propiedades",
                      "• Professional: $999/mes - 500 propiedades",
                      "• Enterprise: desde $1,999/mes - Ilimitado",
                      "Software patentado y protegido"])
    
    # Future Development
    add_content_slide(prs,
                     "Desarrollo Futuro",
                     ["Próximas Características:",
                      "• Realidad Virtual para visitas",
                      "• Blockchain para contratos",
                      "• App móvil nativa",
                      "• Análisis predictivo avanzado"])
    
    # Contact
    add_content_slide(prs,
                     "Contacto",
                     ["Información de Licenciamiento:",
                      "Email: sales@airealestate.com",
                      "Teléfono: (555) 123-4567",
                      "Web: www.airealestate.com",
                      "¡Transforme su agencia inmobiliaria con IA!"])
    
    # Save the presentation
    prs.save('Presentacion_Bienes_Raices_IA.pptx')

if __name__ == '__main__':
    create_presentation()
